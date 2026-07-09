#!/usr/bin/env python3
"""Generate Woodley's Jewelers UX methodology PowerPoint presentation."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "Woodleys-UX-Methodology.pptx"

# Brand palette
IVORY = RGBColor(250, 248, 244)
DARK = RGBColor(44, 40, 37)
GOLD = RGBColor(154, 123, 79)
MUTED = RGBColor(107, 101, 96)
WHITE = RGBColor(255, 255, 255)
CHAMPAGNE = RGBColor(245, 240, 232)


def set_slide_bg(slide, color=IVORY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.08)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(10), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_footer(slide, text="Woodley's Jewelers — UX Methodology — July 2026"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(9), Inches(0.3))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    add_accent_bar(slide, top=Inches(2.8), height=Inches(0.06))

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(8.5), Inches(1.2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.1), Inches(8.5), Inches(1.5))
        stf = sub_box.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(18)
        sp.font.color.rgb = MUTED
        sp.alignment = PP_ALIGN.CENTER

    add_footer(slide)
    return slide


def add_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CHAMPAGNE)
    add_accent_bar(slide, top=Inches(3.5), height=Inches(0.05))

    box = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(8.5), Inches(1))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
        stf = sbox.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(16)
        sp.font.color.rgb = MUTED
        sp.alignment = PP_ALIGN.CENTER

    add_footer(slide)
    return slide


def add_content_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.7))
    ttf = title_box.text_frame
    ttf.text = title
    tp = ttf.paragraphs[0]
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    body_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.15), Inches(8.5), Inches(5.5))
    tf = body_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(16 if level == 0 else 14)
        p.font.color.rgb = DARK if level == 0 else MUTED
        p.space_after = Pt(10)

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.35), Inches(8.5), Inches(0.55))
        ntf = note_box.text_frame
        ntf.text = note
        np = ntf.paragraphs[0]
        np.font.size = Pt(12)
        np.font.italic = True
        np.font.color.rgb = GOLD

    add_footer(slide)
    return slide


def add_table_slide(prs, title, headers, rows, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.7))
    ttf = title_box.text_frame
    ttf.text = title
    tp = ttf.paragraphs[0]
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.4 * (len(rows) + 2))
    )
    table = table_shape.table

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = CHAMPAGNE
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = DARK

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.35), Inches(8.5), Inches(0.55))
        ntf = note_box.text_frame
        ntf.text = note
        np = ntf.paragraphs[0]
        np.font.size = Pt(12)
        np.font.italic = True
        np.font.color.rgb = GOLD

    add_footer(slide)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title
    add_title_slide(
        prs,
        "Woodley's Jewelers",
        "UX Principles, Logic & Design Methodology\nJuly 2026",
    )

    add_content_slide(
        prs,
        "Presentation Purpose",
        [
            "Explain the UX thinking behind the Woodley's Jewelers web application.",
            "Document principles, information architecture, and interaction patterns.",
            "Support stakeholder conversations with designers, developers, and business owners.",
            "Capture decisions that shape both the storefront and admin dashboard.",
        ],
        note="Designed as a luxury showroom experience — not a generic e-commerce template.",
    )

    add_content_slide(
        prs,
        "Executive Summary",
        [
            "Three balanced goals drive every design decision:",
            ("1. Inspire trust and craftsmanship", 1),
            ("Editorial layout, warm typography, appointment-first cues", 2),
            ("2. Support real shopping workflows", 1),
            ("Browse → product detail → cart → checkout", 2),
            ("3. Give staff a focused admin workspace", 1),
            ("Manage inventory and merchandising without clutter", 2),
            "Storefront and dashboard are intentionally separate.",
            "Filters use URL state; destructive actions require explicit confirmation.",
        ],
    )

    add_section_slide(prs, "Product Context", "Who we design for and why it matters")

    add_table_slide(
        prs,
        "Who Uses the App",
        ["Audience", "Primary Goals"],
        [
            ["Shoppers", "Discover jewelry, explore categories, purchase or book a visit"],
            ["Members", "Manage account, orders, and profile"],
            ["Staff / Admins", "Manage products, orders, featured items, inventory"],
        ],
    )

    add_content_slide(
        prs,
        "High-Touch Retail Model",
        [
            "Woodley's is a relationship-driven jeweler — not a volume marketplace.",
            "Phone number and Book Visit are prominent in the header.",
            "Product pages emphasize imagery, craftsmanship, and category context.",
            "Admin tools prioritize inventory control and merchandising over bulk operations.",
            "Service categories (repairs, appraisals, sizing) reflect ongoing customer care.",
        ],
        note="The UX supports online shopping while nudging toward in-store relationships.",
    )

    add_section_slide(prs, "Core UX Principles", "Six principles that guide every surface")

    add_content_slide(
        prs,
        "Principles 1–3",
        [
            "1. Editorial Luxury, Not Marketplace Clutter",
            ("Serif headlines, champagne/ivory surfaces, curated grids", 1),
            "2. Appointment-First Retail",
            ("Visit CTAs, service categories, brand storytelling on homepage", 1),
            "3. Clear Separation of Contexts",
            ("Storefront shell for shoppers; dedicated dashboard for admins", 1),
        ],
    )

    add_content_slide(
        prs,
        "Principles 4–6",
        [
            "4. Reversible Before Permanent",
            ("Deactivate hides products; delete only after deactivation + type DELETE", 1),
            "5. State You Can Share and Return To",
            ("Category, search, and sort live in the URL for bookmarking", 1),
            "6. Soft Overlays, Not Blocking Darkness",
            ("Light blur backdrops preserve spatial context behind modals", 1),
        ],
    )

    add_section_slide(prs, "Information Architecture", "How the site is organized")

    add_content_slide(
        prs,
        "Storefront Structure",
        [
            "Home → Shop / Collections by category",
            ("Engagement & Wedding, Diamonds, Custom Jewelry, Fine Jewelry, Watches, Services", 1),
            "Product detail pages (canonical URLs)",
            "Cart → Checkout",
            "Account (member self-service)",
            "Content (About, Contact, policies)",
        ],
        note="Navigation mirrors how customers think — by occasion and type, not internal SKUs.",
    )

    add_content_slide(
        prs,
        "Admin Structure & Access",
        [
            "Shallow global nav: Dashboard, Orders, Products, Settings",
            "Complex product organization lives inside the Products page",
            "Unauthenticated users → sign in",
            "Non-admins → member account area",
            "Dashboard routes excluded from search indexing",
        ],
        note="Admin is a private workspace; members and admins never share the same nav model.",
    )

    add_section_slide(prs, "Storefront Experience", "Customer-facing journeys")

    add_content_slide(
        prs,
        "Homepage Funnel",
        [
            "1. Hero — rotating highlights with pause/play and keyboard support",
            "2. Brand story — why Woodley's",
            "3. Showroom highlights — featured products",
            "4. New releases — freshness and discovery",
            "5. Featured categories — entry points into the catalog",
        ],
        note="Lead with emotion and craftsmanship, then merchandise, then breadth.",
    )

    add_table_slide(
        prs,
        "Product Detail Page (PDP)",
        ["Element", "UX Decision"],
        [
            ["Layout", "Gallery left; sticky purchase panel on desktop"],
            ["Hierarchy", "Title → price → stock → quantity → add to cart"],
            ["Feedback", "Brief added confirmation after add-to-cart"],
            ["Discovery", "Similar products below the fold"],
            ["Context", "Breadcrumbs back to shop and category"],
        ],
        note="Decision moment is visual and emotional; purchase controls reduce scroll friction.",
    )

    add_content_slide(
        prs,
        "Cart & Checkout",
        [
            "Cart drawer (primary) — opens from header when items exist",
            ("Suggestions column beside bag on larger screens", 1),
            "Full cart page — fallback for deep review or direct links",
            "Checkout handles local catalog and Shopify flows",
            "Mixed carts blocked with clear guidance — never silent failure",
        ],
        note="Drawer keeps shoppers in context; explicit messaging prevents checkout frustration.",
    )

    add_section_slide(prs, "Admin Dashboard", "Staff-facing workflows")

    add_content_slide(
        prs,
        "Dashboard Shell",
        [
            "Icon sidebar on small screens; labels on large screens",
            "Solid sidebar background for clear separation from content",
            "Active nav uses warm amber highlighting (brand CTA color)",
            "Header: brand link, personalized welcome, logout",
        ],
        note="Fast wayfinding without sacrificing content space on tablets.",
    )

    add_content_slide(
        prs,
        "Products Page Layout",
        [
            "Full width: title, search, and sort above the split",
            "25% category panel | 75% product lists",
            ("Expandable categories with subcategories and counts", 1),
            "Active products and deactivated products in separate sections",
            "Category nav inside the page — not in the global sidebar",
        ],
        note="Global nav = app-level tasks. Category tree = product-management task.",
    )

    add_table_slide(
        prs,
        "Filtering, Sorting & Actions",
        ["Control / State", "Behavior"],
        [
            ["Category panel", "URL: ?section=services&collection=appraisals"],
            ["Search", "Filters by product name or handle"],
            ["Sort", "A–Z, Z–A, price, featured, recently added/updated"],
            ["Active products", "Feature, Deactivate, Edit"],
            ["Deactivated products", "Activate (primary), Edit, Delete"],
        ],
        note="One control per job. Delete requires typing DELETE in a confirmation dialog.",
    )

    add_section_slide(prs, "Visual Design System", "Consistent luxury language")

    add_table_slide(
        prs,
        "Color, Type & Buttons",
        ["Element", "Decision"],
        [
            ["Background", "Warm ivory (#faf8f4)"],
            ["Text", "Deep brown-gray"],
            ["Accent", "Warm gold for CTAs and highlights"],
            ["Headings", "Playfair Display (serif)"],
            ["UI / body", "Inter (sans-serif)"],
            ["Primary button", "Amber — Add product, Activate"],
            ["Destructive", "Red outline — same shape as secondary"],
        ],
    )

    add_section_slide(prs, "Interaction & Accessibility", "Patterns that build trust")

    add_content_slide(
        prs,
        "Interaction Patterns",
        [
            "Loading: skeleton layouts, text states, route transition bar",
            "Empty states: guidance + CTAs (never dead ends)",
            "Feedback: success/error banners, inline validation, add-to-cart confirmation",
            "Modals: faint blur backdrop, focus trap, Escape to close",
            "Destructive actions: type-to-confirm for permanent delete",
        ],
    )

    add_table_slide(
        prs,
        "Accessibility & Responsive",
        ["Area", "Approach"],
        [
            ["Landmarks", "main, nav, dialog roles"],
            ["Focus", "Moves to main on route change; trap in modals"],
            ["Screen readers", "Route title announced; accessible price labels"],
            ["Motion", "Respects prefers-reduced-motion"],
            ["Mobile header", "Hamburger + full-screen menu"],
            ["Touch targets", "44×44px minimum (e.g. hamburger)"],
        ],
    )

    add_section_slide(prs, "Performance as UX", "Speed is part of the experience")

    add_table_slide(
        prs,
        "Performance Choices",
        ["Choice", "User Benefit"],
        [
            ["ISR / caching (60s)", "Fast loads; reasonably fresh catalog"],
            ["Prefetch on hover", "Snappier navigation to product pages"],
            ["Lazy-loaded charts", "Dashboard home loads without blocking"],
            ["Cart suggestions on open", "Faster initial page load"],
            ["Client-side admin filter/sort", "Instant response after data loads"],
            ["Font display: swap", "Text visible during font load"],
        ],
        note="Optimized for perceived speed and stable layouts — not metrics alone.",
    )

    add_section_slide(prs, "Key Decisions", "Quick reference for stakeholders")

    add_table_slide(
        prs,
        "Decision Summary (1 of 2)",
        ["Decision", "Why"],
        [
            ["Separate storefront & dashboard", "Clear roles, less confusion"],
            ["Categories inside Products page", "Task-specific without cluttering nav"],
            ["URL-based filters", "Shareable, refresh-safe views"],
            ["Sort replaces collection filter", "Sidebar already handles categories"],
            ["Deactivate vs delete", "Reversible mistakes; delete is explicit"],
        ],
    )

    add_table_slide(
        prs,
        "Decision Summary (2 of 2)",
        ["Decision", "Why"],
        [
            ["Type DELETE to confirm", "Prevents accidental data loss"],
            ["Activate = primary button", "Highlights main recovery action"],
            ["No Feature on deactivated", "Prevents invalid featured state"],
            ["Faint blur modals", "Premium feel; context preserved"],
            ["Appointment CTAs in header", "Matches high-touch retail model"],
        ],
    )

    add_content_slide(
        prs,
        "Stakeholder Talking Points",
        [
            "Business owners:",
            ('"Designed like a showroom, not a warehouse. Staff can hide products safely."', 1),
            "Designers:",
            ('"One visual language: ivory, gold, serif headlines, consistent button hierarchy."', 1),
            "Developers:",
            ('"State in the URL. Intentional layout splits. First-class empty/loading/error states."', 1),
        ],
    )

    add_title_slide(
        prs,
        "Thank You",
        "Woodley's Jewelers — UX Methodology\nQuestions & discussion",
    )

    prs.save(OUTPUT)
    print(f"Presentation written to {OUTPUT}")


if __name__ == "__main__":
    build_presentation()
